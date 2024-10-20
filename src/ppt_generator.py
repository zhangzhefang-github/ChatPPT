import os
from pptx import Presentation
from utils import remove_all_slides
from pathlib import Path

def generate_presentation(powerpoint_data, template_path: str, output_path: str):
    if not os.path.exists(template_path):
        raise FileNotFoundError(f"Template file '{template_path}' does not exist.")

    prs = Presentation(template_path)
    remove_all_slides(prs)
    prs.core_properties.title = powerpoint_data.title

    for slide in powerpoint_data.slides:
        if slide.layout >= len(prs.slide_layouts):
            slide_layout = prs.slide_layouts[0]
        else:
            slide_layout = prs.slide_layouts[slide.layout]

        new_slide = prs.slides.add_slide(slide_layout)

        if new_slide.shapes.title:
            new_slide.shapes.title.text = slide.content.title

        for shape in new_slide.shapes:
            if shape.has_text_frame and not shape == new_slide.shapes.title:
                text_frame = shape.text_frame
                text_frame.clear()
                for point in slide.content.bullet_points:
                    p = text_frame.add_paragraph()
                    p.text = point
                    p.level = 0
                break


        # 提前准备图片路径和占位符
        if slide.content.image_paths:
            image_placeholders = [shape for shape in new_slide.placeholders if shape.placeholder_format.type == 18]
            
            # 使用 Path 来处理路径，并且只保留存在的图片路径
            existing_images = [Path(image_path) for image_path in slide.content.image_paths if Path(image_path).exists()]

            # 检查图片数量和占位符数量是否匹配
            if len(existing_images) > len(image_placeholders):
                print(f"Warning: There are more images ({len(existing_images)}) than placeholders ({len(image_placeholders)}).")
            elif len(existing_images) < len(image_placeholders):
                print(f"Warning: There are fewer images ({len(existing_images)}) than placeholders ({len(image_placeholders)}).")

            # 同时遍历现有图片和占位符，并插入图片
            for image_path, placeholder in zip(existing_images, image_placeholders):
                try:
                    new_slide.shapes.add_picture(str(image_path), placeholder.left, placeholder.top, placeholder.width, placeholder.height)
                except Exception as e:
                    print(f"Error inserting image {image_path}: {e}")

    prs.save(output_path)
    print(f"Presentation saved to '{output_path}'")
